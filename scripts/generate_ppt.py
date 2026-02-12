#!/usr/bin/env python3
"""
PPT Generator - python-pptx 기반 와인 테이스팅 노트 슬라이드 생성
Usage: python generate_ppt.py <input.json> <output.pptx>

입력 JSON 형식:
{
  "slides": [{ ... wine data ... }],
  "logoPath": "/tmp/logo.jpg",   (optional)
  "iconPath": "/tmp/icon.jpg"    (optional)
}
"""

import sys
import os
import json
import subprocess

# 런타임 패키지 설치 (Vercel Lambda 등 패키지 미설치 환경)
_PKG_DIR = '/tmp/python_packages'

def _ensure_packages():
    """python-pptx, Pillow가 없으면 /tmp에 자동 설치"""
    try:
        import pptx  # noqa: F401
        return
    except ImportError:
        pass
    os.makedirs(_PKG_DIR, exist_ok=True)
    sys.path.insert(0, _PKG_DIR)
    subprocess.check_call([
        sys.executable, '-m', 'pip', 'install',
        'python-pptx', 'Pillow',
        '-t', _PKG_DIR, '--quiet', '--disable-pip-version-check'
    ], stdout=subprocess.DEVNULL)

_ensure_packages()

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════
# 와인 테마 컬러 팔레트
# ═══════════════════════════════════════════════════
COLORS = {
    'BG_CREAM':        RGBColor(0xFA, 0xF7, 0xF2),
    'BG_BOTTLE_AREA':  RGBColor(0xF5, 0xF0, 0xEA),
    'BURGUNDY':        RGBColor(0x72, 0x2F, 0x37),
    'BURGUNDY_DARK':   RGBColor(0x5A, 0x25, 0x2C),
    'BURGUNDY_LIGHT':  RGBColor(0xF2, 0xE8, 0xEA),
    'GOLD':            RGBColor(0xB8, 0x97, 0x6A),
    'GOLD_LIGHT':      RGBColor(0xD4, 0xC4, 0xA8),
    'TEXT_PRIMARY':    RGBColor(0x2C, 0x2C, 0x2C),
    'TEXT_SECONDARY':  RGBColor(0x5A, 0x5A, 0x5A),
    'TEXT_MUTED':      RGBColor(0x8A, 0x8A, 0x8A),
    'TEXT_ON_DARK':    RGBColor(0xFF, 0xFF, 0xFF),
    'CARD_BORDER':     RGBColor(0xE0, 0xD5, 0xC8),
    'DIVIDER':         RGBColor(0xD4, 0xC4, 0xA8),
    'DIVIDER_LIGHT':   RGBColor(0xE8, 0xDD, 0xD0),
    'WHITE':           RGBColor(0xFF, 0xFF, 0xFF),
}

FONT_MAIN = 'Gowun Dodum'
FONT_EN = 'Georgia'        # 영문 세리프 (와인명 영문, 라벨, 와이너리)

# 슬라이드 크기 (인치) - 세로 A4
SLIDE_W = 7.5
SLIDE_H = 10.0


def inches(val):
    return Inches(val)


def add_line(slide, x, y, w, color, width_pt=0.75):
    """얇은 수평선 추가"""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        inches(x), inches(y),
        inches(x + w), inches(y)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    return connector


def add_textbox(slide, x, y, w, h, text='', font_size=9, color=None,
                bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                v_anchor=MSO_ANCHOR.TOP, font_name=FONT_MAIN,
                line_spacing=None, word_wrap=True):
    """텍스트박스 추가 헬퍼"""
    txBox = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color or COLORS['TEXT_PRIMARY']
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment

    if line_spacing:
        p.line_spacing = Pt(line_spacing)

    tf.auto_size = None

    # vertical anchor 설정
    from pptx.oxml.ns import qn
    bodyPr = tf.paragraphs[0]._p.getparent().getparent().find(qn('a:bodyPr'))
    if bodyPr is not None:
        anchor_map = {
            MSO_ANCHOR.TOP: 't',
            MSO_ANCHOR.MIDDLE: 'ctr',
            MSO_ANCHOR.BOTTOM: 'b',
        }
        bodyPr.set('anchor', anchor_map.get(v_anchor, 't'))

    return txBox


def _set_fill_transparency(shape, transparency):
    """shape fill에 투명도 적용 (0~100, 100=완전 투명)"""
    if transparency <= 0:
        return
    from pptx.oxml.ns import qn
    from lxml import etree
    # spPr → solidFill → srgbClr 찾기
    spPr = shape._element.find(qn('a:solidFill'), shape._element.nsmap)
    if spPr is None:
        spPr_parent = shape._element.find('.//' + qn('a:solidFill'))
        if spPr_parent is not None:
            srgbClr = spPr_parent.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha.set('val', str(int((100 - transparency) * 1000)))
                return
    # fallback: XML에서 직접 찾기
    for srgbClr in shape._element.iter(qn('a:srgbClr')):
        alpha = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha.set('val', str(int((100 - transparency) * 1000)))
        break


def add_rounded_rect(slide, x, y, w, h, fill_color=None, border_color=None,
                     border_width=0.5, transparency=0):
    """둥근 사각형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        inches(x), inches(y), inches(w), inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if transparency > 0:
            _set_fill_transparency(shape, transparency)
    else:
        shape.fill.background()

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()

    return shape


def add_rect(slide, x, y, w, h, fill_color=None, transparency=0):
    """일반 사각형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        inches(x), inches(y), inches(w), inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if transparency > 0:
            _set_fill_transparency(shape, transparency)
    else:
        shape.fill.background()

    shape.line.fill.background()
    return shape


def add_label_badge(slide, text, x, y, w, h=0.22):
    """버건디 라벨 뱃지"""
    shape = add_rounded_rect(slide, x, y, w, h,
                             fill_color=COLORS['BURGUNDY'],
                             border_color=None)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8.5)
    p.font.name = FONT_MAIN
    p.font.color.rgb = COLORS['TEXT_ON_DARK']
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    from pptx.oxml.ns import qn
    bodyPr = tf.paragraphs[0]._p.getparent().getparent().find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')

    return shape


def crop_whitespace(img):
    """이미지 여백(투명/흰색) 제거"""
    from PIL import ImageOps

    if img.mode != 'RGBA':
        img = img.convert('RGBA')

    w, h = img.size
    alpha = img.split()[3]
    amin, _ = alpha.getextrema()

    # 투명 배경 이미지: 알파 채널 기준 크롭
    if amin < 128:
        bbox = alpha.getbbox()
        if bbox:
            cw, ch = bbox[2] - bbox[0], bbox[3] - bbox[1]
            if cw < w * 0.95 or ch < h * 0.95:
                pad = max(3, min(w, h) // 80)
                bbox = (max(0, bbox[0] - pad), max(0, bbox[1] - pad),
                        min(w, bbox[2] + pad), min(h, bbox[3] + pad))
                return img.crop(bbox)
        return img

    # 불투명 이미지: 밝기 기준 크롭
    gray = img.convert('L')
    inv = ImageOps.invert(gray)
    bw = inv.point(lambda x: 255 if x > 25 else 0)
    bbox = bw.getbbox()
    if bbox:
        cw, ch = bbox[2] - bbox[0], bbox[3] - bbox[1]
        if cw < w * 0.95 or ch < h * 0.95:
            pad = max(3, min(w, h) // 80)
            bbox = (max(0, bbox[0] - pad), max(0, bbox[1] - pad),
                    min(w, bbox[2] + pad), min(h, bbox[3] + pad))
            return img.crop(bbox)

    return img


def embed_font_in_pptx(pptx_path):
    """PPTX에 Gowun Dodum 폰트 임베딩 (ECMA-376 §13.1)"""
    import zipfile
    import uuid
    import tempfile
    import shutil

    # 폰트 파일 찾기
    font_candidates = [
        os.path.join(os.path.dirname(os.path.dirname(__file__)), 'data', 'fonts', 'GowunDodum-Regular.ttf'),
        os.path.join(os.path.expanduser('~'), 'AppData', 'Local', 'Microsoft', 'Windows', 'Fonts', 'GowunDodum-Regular.ttf'),
        os.path.join('C:\\', 'Windows', 'Fonts', 'GowunDodum-Regular.ttf'),
    ]
    font_path = None
    for fp in font_candidates:
        if os.path.exists(fp):
            font_path = fp
            break

    if not font_path:
        print(f"Warning: Gowun Dodum font not found, skipping embed", file=sys.stderr)
        return

    with open(font_path, 'rb') as f:
        font_data = bytearray(f.read())

    # GUID 생성 및 난독화
    guid = uuid.uuid4()
    guid_hex = guid.hex  # 32 hex chars
    key = bytes(int(guid_hex[(15 - i) * 2:(15 - i) * 2 + 2], 16) for i in range(16))

    obfuscated = bytearray(font_data)
    for i in range(32):
        obfuscated[i] ^= key[i % 16]

    guid_str = str(guid).upper()
    font_filename = f'{{{guid_str}}}.fntdata'
    r_id = 'rIdEmbedFont1'

    # PPTX (ZIP) 수정
    tmp_path = pptx_path + '.tmp'
    with zipfile.ZipFile(pptx_path, 'r') as zin, zipfile.ZipFile(tmp_path, 'w') as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)

            if item.filename == '[Content_Types].xml':
                text = data.decode('utf-8')
                text = text.replace('</Types>',
                    f'<Override PartName="/ppt/fonts/{font_filename}" ContentType="application/x-fontdata"/></Types>')
                data = text.encode('utf-8')

            elif item.filename == 'ppt/_rels/presentation.xml.rels':
                text = data.decode('utf-8')
                text = text.replace('</Relationships>',
                    f'<Relationship Id="{r_id}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/font" Target="fonts/{font_filename}"/></Relationships>')
                data = text.encode('utf-8')

            elif item.filename == 'ppt/presentation.xml':
                text = data.decode('utf-8')
                import re as _re
                embed_xml = f'<p:embeddedFontLst><p:embeddedFont><p:font typeface="Gowun Dodum"/><p:regular r:id="{r_id}"/></p:embeddedFont></p:embeddedFontLst>'
                # OOXML 스키마 순서: notesSz → embeddedFontLst → defaultTextStyle
                # notesSz는 self-closing (<p:notesSz .../>) 또는 closing (</p:notesSz>) 둘 다 처리
                if _re.search(r'<p:notesSz[^>]*/>', text):
                    text = _re.sub(r'(<p:notesSz[^>]*/>)', r'\1' + embed_xml, text)
                elif '</p:notesSz>' in text:
                    text = text.replace('</p:notesSz>', f'</p:notesSz>{embed_xml}')
                elif '<p:defaultTextStyle' in text:
                    text = text.replace('<p:defaultTextStyle', f'{embed_xml}<p:defaultTextStyle')
                else:
                    text = text.replace('</p:presentation>', f'{embed_xml}</p:presentation>')
                data = text.encode('utf-8')

            zout.writestr(item, data)

        # 폰트 파일 추가
        zout.writestr(f'ppt/fonts/{font_filename}', bytes(obfuscated))

    shutil.move(tmp_path, pptx_path)


def add_tasting_note_slide(prs, data, logo_path=None, icon_path=None):
    """단일 와인 테이스팅 노트 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # ════════════════════════════════════════════
    # 1. 배경: 흰색
    # ════════════════════════════════════════════
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['WHITE']

    # 2. 좌측 병 영역 배경 패널
    add_rect(slide, 0, 0.90, 2.10, 8.10,
             fill_color=COLORS['BG_BOTTLE_AREA'], transparency=40)

    # ════════════════════════════════════════════
    # 3. HEADER 영역 - Logo
    # ════════════════════════════════════════════
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path,
                                     inches(0.20), inches(0.20),
                                     inches(1.49), inches(0.57))
        except Exception:
            pass

    # 와이너리 태그라인
    winery_desc = data.get('wineryDescription', '')
    if winery_desc:
        tagline = winery_desc.split('.')[0].strip()
        if not tagline:
            tagline = winery_desc.split('。')[0].strip()
        if tagline:
            add_textbox(slide, 1.76, 0.20, 5.20, 0.24,
                        text=tagline, font_size=9,
                        color=COLORS['TEXT_MUTED'], italic=True,
                        font_name=FONT_EN)

    # 4-5. 헤더 구분선 (버건디 + 골드 이중선)
    add_line(slide, 0.20, 0.84, 7.10, COLORS['BURGUNDY'], 1.0)
    add_line(slide, 0.20, 0.87, 7.10, COLORS['GOLD_LIGHT'], 0.75)

    # ════════════════════════════════════════════
    # 6. 와인명 카드 배경 (둥근 사각형)
    # ════════════════════════════════════════════
    add_rounded_rect(slide, 2.05, 0.97, 5.20, 0.76,
                     fill_color=COLORS['BURGUNDY_LIGHT'],
                     border_color=COLORS['CARD_BORDER'],
                     transparency=20)

    # 7. 와인명 텍스트 (한글 + 영문)
    name_kr = data.get('nameKr', '')
    name_en = data.get('nameEn', '')

    # 한글명 앞 영어약어 2자 제거
    import re
    name_kr_clean = re.sub(r'^[A-Za-z]{2}\s+', '', name_kr)

    txBox = slide.shapes.add_textbox(inches(2.20), inches(1.00), inches(4.90), inches(0.72))
    tf = txBox.text_frame
    tf.word_wrap = True

    # 한글명 paragraph
    p_kr = tf.paragraphs[0]
    p_kr.text = name_kr_clean
    p_kr.font.size = Pt(14.5)
    p_kr.font.name = FONT_MAIN
    p_kr.font.color.rgb = COLORS['BURGUNDY_DARK']
    p_kr.font.bold = True
    p_kr.alignment = PP_ALIGN.LEFT

    # 영문명 paragraph (Georgia Italic)
    if name_en:
        p_en = tf.add_paragraph()
        p_en.text = name_en
        p_en.font.size = Pt(10.5)
        p_en.font.name = FONT_EN
        p_en.font.color.rgb = COLORS['TEXT_SECONDARY']
        p_en.font.bold = False
        p_en.font.italic = True
        p_en.alignment = PP_ALIGN.LEFT
        p_en.space_before = Pt(2)

    # 8. 와인명 하단 구분선
    add_line(slide, 2.20, 1.82, 4.90, COLORS['DIVIDER'], 0.75)

    # ════════════════════════════════════════════
    # 9-10. 지역
    # ════════════════════════════════════════════
    add_label_badge(slide, '지역', 2.12, 1.97, 0.55)

    region = data.get('region', '')
    country_en = data.get('countryEn', '') or data.get('country', '')
    region_text = f"{country_en}, {region}" if region else (country_en or '-')
    add_textbox(slide, 2.75, 1.96, 4.40, 0.24,
                text=region_text, font_size=9.5)

    # 구분선
    add_line(slide, 2.20, 2.35, 4.90, COLORS['DIVIDER_LIGHT'], 0.5)

    # 11-12. 품종
    add_label_badge(slide, '품종', 2.12, 2.42, 0.55)
    add_textbox(slide, 2.75, 2.41, 4.40, 0.40,
                text=data.get('grapeVarieties', '') or '-',
                font_size=9.5)

    # ════════════════════════════════════════════
    # 13-14. 빈티지
    # ════════════════════════════════════════════
    add_label_badge(slide, '빈티지', 2.12, 3.02, 0.65)

    vintage = data.get('vintage', '') or '-'
    add_textbox(slide, 2.85, 2.98, 0.75, 0.28,
                text=vintage, font_size=13,
                color=COLORS['BURGUNDY'], bold=True)

    vintage_note = data.get('vintageNote', '')
    if vintage_note:
        add_textbox(slide, 3.65, 3.00, 3.50, 0.40,
                    text=vintage_note, font_size=8,
                    color=COLORS['TEXT_SECONDARY'])

    # ════════════════════════════════════════════
    # 15-16. 양조
    # ════════════════════════════════════════════
    add_label_badge(slide, '양조', 2.12, 3.62, 0.55)

    winemaking_text = data.get('winemaking', '') or '-'
    alcohol = data.get('alcoholPercentage', '')
    if alcohol:
        winemaking_text += f"\n알코올: {alcohol}"
    add_textbox(slide, 2.15, 3.90, 5.00, 1.23,
                text=winemaking_text, font_size=9,
                line_spacing=12)

    # ════════════════════════════════════════════
    # 17-18. 테이스팅 노트 (핵심 영역)
    # ════════════════════════════════════════════
    add_rounded_rect(slide, 2.05, 5.30, 5.20, 2.72,
                     fill_color=COLORS['BURGUNDY_LIGHT'],
                     border_color=COLORS['CARD_BORDER'],
                     transparency=30)

    # TASTING NOTE 라벨
    add_label_badge(slide, 'TASTING NOTE', 2.12, 5.35, 1.32, h=0.22)

    # 테이스팅 노트 내용
    tasting_items = [
        ('Color', data.get('colorNote', '')),
        ('Nose', data.get('noseNote', '')),
        ('Palate', data.get('palateNote', '')),
        ('Potential', data.get('agingPotential', '')),
        # ('Serving', data.get('servingTemp', '')),  # 제외
    ]

    txBox = slide.shapes.add_textbox(inches(2.15), inches(5.62), inches(5.00), inches(2.32))
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for label, value in tasting_items:
        if not value:
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(6)

        # Label run (Georgia Italic, burgundy)
        run_label = p.add_run()
        run_label.text = label
        run_label.font.size = Pt(8.5)
        run_label.font.name = FONT_EN
        run_label.font.color.rgb = COLORS['BURGUNDY']
        run_label.font.bold = False
        run_label.font.italic = True

        # Value run
        run_value = p.add_run()
        run_value.text = f"\n{value}"
        run_value.font.size = Pt(9)
        run_value.font.name = FONT_MAIN
        run_value.font.color.rgb = COLORS['TEXT_PRIMARY']

    if first:
        # 모든 tasting items가 비어있는 경우
        p = tf.paragraphs[0]
        p.text = '-'
        p.font.size = Pt(9)
        p.font.name = FONT_MAIN
        p.font.color.rgb = COLORS['TEXT_MUTED']

    # from pptx.oxml.ns import qn
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 't')
        # line spacing
        bodyPr.set('lIns', '0')
        bodyPr.set('tIns', '0')
        bodyPr.set('rIns', '0')
        bodyPr.set('bIns', '0')

    # ════════════════════════════════════════════
    # 19. 푸드 페어링
    # ════════════════════════════════════════════
    add_label_badge(slide, '푸드 페어링', 2.12, 8.18, 0.95)
    add_textbox(slide, 2.15, 8.42, 5.00, 0.44,
                text=data.get('foodPairing', '') or '-',
                font_size=9, line_spacing=12)

    # ════════════════════════════════════════════
    # 20. 수상내역 영역
    # ════════════════════════════════════════════
    add_line(slide, 0.15, 9.04, 7.20, COLORS['GOLD_LIGHT'], 0.5)

    awards = data.get('awards', '')
    if awards and awards != 'N/A':
        # 수상 아이콘
        if icon_path and os.path.exists(icon_path):
            try:
                slide.shapes.add_picture(icon_path,
                                         inches(0.25), inches(9.08),
                                         inches(0.22), inches(0.28))
            except Exception:
                pass

        txBox = slide.shapes.add_textbox(inches(0.52), inches(9.07), inches(6.70), inches(0.30))
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        run_label = p.add_run()
        run_label.text = 'AWARDS  '
        run_label.font.size = Pt(8)
        run_label.font.name = FONT_MAIN
        run_label.font.color.rgb = COLORS['GOLD']
        run_label.font.bold = True

        run_value = p.add_run()
        run_value.text = awards
        run_value.font.size = Pt(9)
        run_value.font.name = FONT_MAIN
        run_value.font.color.rgb = COLORS['TEXT_PRIMARY']

        from pptx.oxml.ns import qn
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')

    # ════════════════════════════════════════════
    # 21-23. FOOTER 영역
    # ════════════════════════════════════════════
    add_line(slide, 0.20, 9.52, 7.10, COLORS['BURGUNDY'], 2.0)
    add_line(slide, 0.20, 9.55, 7.10, COLORS['GOLD_LIGHT'], 0.75)

    # Footer 로고
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path,
                                     inches(0.09), inches(9.68),
                                     inches(0.95), inches(0.25))
        except Exception:
            pass

    # 회사 정보
    add_textbox(slide, 1.12, 9.69, 2.76, 0.24,
                text='T. 02-786-3136  |  www.cavedevin.co.kr',
                font_size=7, color=COLORS['TEXT_MUTED'],
                alignment=PP_ALIGN.RIGHT)

    # ════════════════════════════════════════════
    # 24. 병 이미지 (있을 때만)
    # ════════════════════════════════════════════
    bottle_path = data.get('bottleImagePath', '')
    if bottle_path and os.path.exists(bottle_path):
        try:
            from PIL import Image
            img = Image.open(bottle_path)

            # 팔레트(P) 모드 이미지 → RGBA로 변환
            if img.mode in ('P', 'PA'):
                img = img.convert('RGBA')
            elif img.mode in ('L', 'RGB'):
                img = img.convert('RGBA')

            # 여백 크롭
            img = crop_whitespace(img)

            converted_path = bottle_path + '_converted.png'
            img.save(converted_path, 'PNG')

            # 비율 보정하여 영역에 맞추기
            area_w = 1.60
            area_h = 5.50
            area_x = 0.25
            area_y = 2.00

            iw, ih = img.size
            img.close()

            ratio = iw / ih
            ar = area_w / area_h

            if ratio > ar:
                fw = area_w
                fh = area_w / ratio
            else:
                fh = area_h
                fw = area_h * ratio

            ox = area_x + (area_w - fw) / 2
            oy = area_y + (area_h - fh) * 0.3

            slide.shapes.add_picture(converted_path,
                                     inches(ox), inches(oy),
                                     inches(fw), inches(fh))

            try:
                os.remove(converted_path)
            except Exception:
                pass

        except Exception as e:
            print(f"Warning: Failed to add bottle image: {e}", file=sys.stderr)


def main():
    if len(sys.argv) < 3:
        print("Usage: python generate_ppt.py <input.json> <output.pptx>", file=sys.stderr)
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2]

    # JSON 입력 읽기
    with open(input_path, 'r', encoding='utf-8') as f:
        payload = json.load(f)

    slides_data = payload.get('slides', [])
    logo_path = payload.get('logoPath', '')
    icon_path = payload.get('iconPath', '')

    if not slides_data:
        print("Error: No slides data", file=sys.stderr)
        sys.exit(1)

    # Presentation 생성
    prs = Presentation()

    # 세로 A4 레이아웃 (7.5 x 10 인치)
    prs.slide_width = inches(SLIDE_W)
    prs.slide_height = inches(SLIDE_H)

    # 슬라이드 생성
    for data in slides_data:
        add_tasting_note_slide(prs, data, logo_path, icon_path)

    # 저장
    prs.save(output_path)
    print(f"Saved: {output_path} ({len(slides_data)} slides)", file=sys.stderr)

    # 폰트 임베딩
    try:
        embed_font_in_pptx(output_path)
        print("Font embedded: Gowun Dodum", file=sys.stderr)
    except Exception as e:
        print(f"Font embedding warning: {e}", file=sys.stderr)

    # 검증: 재오픈 테스트
    try:
        verify_prs = Presentation(output_path)
        assert len(verify_prs.slides) == len(slides_data), \
            f"Slide count mismatch: expected {len(slides_data)}, got {len(verify_prs.slides)}"
        print("Verification passed", file=sys.stderr)
    except Exception as e:
        print(f"Verification warning: {e}", file=sys.stderr)
        # 검증 실패해도 원본 파일은 유지

    print("OK", file=sys.stdout)


if __name__ == '__main__':
    main()
